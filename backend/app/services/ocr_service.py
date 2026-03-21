import os
from typing import Optional
from PIL import Image
import pytesseract


def extract_text_from_image(image_path: str, language: str = "eng") -> str:
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image, lang=language)
        return text.strip()
    except Exception as e:
        return f"OCR Error: {str(e)}"


def extract_text_from_pdf(pdf_path: str, language: str = "eng") -> str:
    try:
        from pdf2image import convert_from_path

        images = convert_from_path(pdf_path)
        all_text = []
        for i, image in enumerate(images):
            page_text = pytesseract.image_to_string(image, lang=language)
            all_text.append(f"--- Page {i + 1} ---\n{page_text}")
        return "\n\n".join(all_text)
    except Exception as e:
        return f"PDF OCR Error: {str(e)}"


def extract_text_from_file(
    file_path: str, file_extension: str, language: str = "eng"
) -> str:
    ext = file_extension.lower()

    if ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".tif"}:
        return extract_text_from_image(file_path, language)
    elif ext == ".pdf":
        return extract_text_from_pdf(file_path, language)
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    elif ext in {".doc", ".docx"}:
        return extract_text_from_docx(file_path)
    elif ext in {".xls", ".xlsx"}:
        return extract_text_from_excel(file_path)
    elif ext in {".ppt", ".pptx"}:
        return extract_text_from_pptx(file_path)
    elif ext == ".csv":
        return extract_text_from_csv(file_path)
    else:
        return extract_raw_text(file_path)


def extract_text_from_docx(file_path: str) -> str:
    try:
        from docx import Document

        doc = Document(file_path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        return f"Docx Error: {str(e)}"


def extract_text_from_excel(file_path: str) -> str:
    try:
        import openpyxl

        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_parts = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            text_parts.append(f"Sheet: {sheet}")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    [str(cell) if cell is not None else "" for cell in row]
                )
                if row_text.strip():
                    text_parts.append(row_text)
        return "\n".join(text_parts)
    except Exception as e:
        return f"Excel Error: {str(e)}"


def extract_text_from_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text_parts.append(f"Slide {i + 1}: " + " ".join(slide_text))
        return "\n".join(text_parts)
    except Exception as e:
        return f"PPTX Error: {str(e)}"


def extract_text_from_csv(file_path: str) -> str:
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"CSV Error: {str(e)}"


def extract_raw_text(file_path: str) -> str:
    try:
        with open(file_path, "rb") as f:
            content = f.read()
            return content.decode("utf-8", errors="ignore")[:10000]
    except Exception as e:
        return f"Read Error: {str(e)}"
